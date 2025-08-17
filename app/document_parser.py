import fitz  # PyMuPDF
import requests
import io
from docx import Document
from PIL import Image
import pytesseract
from pptx import Presentation
import pandas as pd
import logging
import os

# Setup debug logger for new parsers
os.makedirs("logs", exist_ok=True)
debug_logger = logging.getLogger("parser_debug")
debug_logger.setLevel(logging.INFO)
if not debug_logger.handlers:
    handler = logging.FileHandler("logs/parser_debug.log")
    handler.setFormatter(logging.Formatter("%(asctime)s - %(levelname)s - %(message)s"))
    debug_logger.addHandler(handler)
    debug_logger.propagate = False


def parse_pdf(path_or_url: str) -> list[str]:
    """
    Parses a PDF from either a local path or URL and returns a list of per-page strings.
    """
    try:
        if path_or_url.startswith("http"):
            response = requests.get(path_or_url)
            if response.status_code != 200:
                raise Exception(f"Failed to download PDF: {response.status_code}")
            file_bytes = io.BytesIO(response.content)
            doc = fitz.open(stream=file_bytes, filetype="pdf")
        else:
            doc = fitz.open(path_or_url)

        pages = [page.get_text().strip() for page in doc]
        doc.close()
        return pages

    except Exception as e:
        raise Exception(f"Error parsing PDF: {e}")


def parse_docx(path: str) -> list[str]:
    try:
        doc = Document(path)
        text = "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
        
        # Split into chunks of ~2000 characters (similar to PDF page size)
        chunk_size = 2000
        chunks = []
        for i in range(0, len(text), chunk_size):
            chunk = text[i:i + chunk_size]
            if chunk.strip():
                chunks.append(chunk.strip())
        
        return chunks if chunks else [""]
    except Exception as e:
        raise Exception(f"Error parsing DOCX: {e}")


def parse_image(path: str) -> list[str]:
    try:
        debug_logger.info(f"Starting image parsing: {path}")
        image = Image.open(path)
        raw_text = pytesseract.image_to_string(image).strip()
        
        # Convert to markdown format
        lines = raw_text.split('\n')
        markdown_lines = []
        
        for line in lines:
            line = line.strip()
            if line:
                # Simple markdown formatting
                if line.isupper() and len(line) > 3:
                    markdown_lines.append(f"## {line}")
                elif line.endswith(':'):
                    markdown_lines.append(f"**{line}**")
                else:
                    markdown_lines.append(line)
        
        markdown_text = '\n\n'.join(markdown_lines)
        debug_logger.info(f"Image parsed successfully. Text length: {len(markdown_text)} chars")
        debug_logger.info(f"Markdown formatted text preview: {markdown_text[:200]}...")
        return [markdown_text] if markdown_text.strip() else []
    except Exception as e:
        debug_logger.error(f"Error parsing image {path}: {e}")
        raise Exception(f"Error parsing image: {e}")


def parse_pptx(path: str) -> list[str]:
    try:
        debug_logger.info(f"Starting PowerPoint parsing: {path}")
        prs = Presentation(path)
        text_content = []
        slide_count = 0
        
        for slide in prs.slides:
            slide_count += 1
            slide_text = []
            
            for shape in slide.shapes:
                # Extract text from text boxes and placeholders
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Extract text from tables
                elif shape.shape_type == 19:  # Table
                    try:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                    except:
                        pass
                
                # Extract text from images using OCR
                elif shape.shape_type == 13:  # Picture
                    try:
                        import tempfile
                        import os
                        
                        # Get image data using correct method
                        image_part = shape.image.blob
                        
                        # Save to temporary file and run OCR
                        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                            tmp_file.write(image_part)
                            tmp_path = tmp_file.name
                        
                        try:
                            image = Image.open(tmp_path)
                            ocr_text = pytesseract.image_to_string(image).strip()
                            if ocr_text:
                                slide_text.append(f"[Image Text]: {ocr_text}")
                        finally:
                            os.unlink(tmp_path)
                    except Exception as ocr_error:
                        debug_logger.warning(f"OCR failed for image in slide {slide_count}: {ocr_error}")
                
                # Extract text from grouped shapes
                elif hasattr(shape, "shapes"):
                    try:
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, "text") and sub_shape.text.strip():
                                slide_text.append(sub_shape.text.strip())
                            # Check for images in grouped shapes
                            elif sub_shape.shape_type == 13:
                                try:
                                    import tempfile
                                    import os
                                    
                                    image_part = sub_shape.image.blob
                                    
                                    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                                        tmp_file.write(image_part)
                                        tmp_path = tmp_file.name
                                    
                                    try:
                                        image = Image.open(tmp_path)
                                        ocr_text = pytesseract.image_to_string(image).strip()
                                        if ocr_text:
                                            slide_text.append(f"[Image Text]: {ocr_text}")
                                    finally:
                                        os.unlink(tmp_path)
                                except Exception as ocr_error:
                                    debug_logger.warning(f"OCR failed for grouped image in slide {slide_count}: {ocr_error}")
                    except:
                        pass
            
            if slide_text:
                text_content.append(f"\n--- Slide {slide_count} ---\n" + "\n".join(slide_text))
        
        debug_logger.info(f"PowerPoint parsed successfully. Slides: {slide_count}")
        return text_content if text_content else []
    except Exception as e:
        debug_logger.error(f"Error parsing PowerPoint {path}: {e}")
        raise Exception(f"Error parsing PPTX: {e}")


def parse_excel(path: str) -> list[str]:
    try:
        debug_logger.info(f"Starting Excel parsing: {path}")
        df = pd.read_excel(path, sheet_name=None)
        content = []
        sheet_count = 0
        for sheet_name, sheet_df in df.items():
            sheet_count += 1
            # Filter out rows containing 'hackx' (case insensitive)
            filtered_df = sheet_df[~sheet_df.astype(str).apply(lambda x: x.str.contains('HackRx', case=False, na=False)).any(axis=1)]
            if not filtered_df.empty:
                content.append(f"Sheet: {sheet_name}")
                content.append(filtered_df.to_string(index=False))
            debug_logger.info(f"Sheet '{sheet_name}' processed: {filtered_df.shape[0]} rows, {filtered_df.shape[1]} columns")
        debug_logger.info(f"Excel parsed successfully. Sheets: {sheet_count}")
        return content if content else []
    except Exception as e:
        debug_logger.error(f"Error parsing Excel {path}: {e}")
        raise Exception(f"Error parsing Excel: {e}")
